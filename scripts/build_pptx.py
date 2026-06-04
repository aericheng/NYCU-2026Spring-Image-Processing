"""Build the 5-minute class-presentation deck (report/Slides.pptx) from the
slide plan in report/PPT_outline.md.

Design: 16:9, consistent title bar + footer, "字少圖大" (few words, big figures).
Each figure is fitted to its box by aspect ratio (never stretched). Slide 2's
15-image contact sheet is generated on the fly from Images/. The 5-minute speaker
timing from PPT_outline.md is written into each slide's notes for the 06/12 talk.

Run with the deblur env python:
  & "C:\\ProgramData\\miniconda3\\envs\\deblur\\python.exe" scripts\\build_pptx.py
"""
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

PROJ = Path(__file__).resolve().parent.parent
FIG = PROJ / "report" / "figures"
IMAGES = PROJ / "Images"
OUT = PROJ / "report" / "Slides.pptx"
GRID = PROJ / "results" / "_ppt" / "images_grid.jpg"   # scratch (results/ is git-ignored)

# Palette (matches the report's print CSS in md_to_pdf.py).
RED = RGBColor(0xC4, 0x4E, 0x52)
BLUE = RGBColor(0x2A, 0x4D, 0x80)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF7, 0xF9, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Microsoft JhengHei"

EMU_IN = 914400
SW, SH = 13.333, 7.5            # 16:9 widescreen, inches


def set_font(run, size, *, bold=False, color=DARK, name=FONT):
    """Set Latin + East-Asian + complex-script typeface so CJK renders in JhengHei."""
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, x, y, w, h, color, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def title_bar(slide, text, page):
    """Top accent bar + title; bottom footer with page number."""
    rect(slide, 0, 0, 0.28, SH, RED)                       # left red spine
    tf = textbox(slide, 0.55, 0.28, SW - 1.1, 0.95, MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    run = p.add_run(); run.text = text
    set_font(run, 28, bold=True, color=BLUE)
    rect(slide, 0.55, 1.22, SW - 1.1, 0.035, RED)          # underline
    # footer
    ftf = textbox(slide, 0.55, SH - 0.5, SW - 1.1, 0.36, MSO_ANCHOR.MIDDLE)
    fp = ftf.paragraphs[0]
    r1 = fp.add_run(); r1.text = "夜間運動模糊的人臉修復 · 影像處理 Term Project"
    set_font(r1, 9, color=GRAY)
    pg = textbox(slide, SW - 1.3, SH - 0.5, 0.75, 0.36, MSO_ANCHOR.MIDDLE)
    pgp = pg.paragraphs[0]; pgp.alignment = PP_ALIGN.RIGHT
    r2 = pgp.add_run(); r2.text = str(page)
    set_font(r2, 11, bold=True, color=RED)


def bullets(slide, items, x, y, w, h, size=17):
    """items: list of str, or (str, level)."""
    tf = textbox(slide, x, y, w, h, MSO_ANCHOR.TOP)
    first = True
    for it in items:
        text, level = (it if isinstance(it, tuple) else (it, 0))
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(10 if level == 0 else 5)
        p.line_spacing = 1.15
        run = p.add_run()
        run.text = ("• " if level == 0 else "– ") + text
        set_font(run, size if level == 0 else size - 2,
                 bold=False, color=DARK if level == 0 else GRAY)
    return tf


def fit_image(slide, path, x, y, max_w, max_h, *, caption=None):
    """Place image fitted (aspect-preserving) and centered inside the box."""
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(max_w / iw, max_h / ih)
    nw, nh = iw * scale, ih * scale
    cx = x + (max_w - nw) / 2
    cy = y + (max_h - nh) / 2
    slide.shapes.add_picture(str(path), Inches(cx), Inches(cy), Inches(nw), Inches(nh))
    if caption:
        ctf = textbox(slide, x, y + max_h - 0.02, max_w, 0.3, MSO_ANCHOR.TOP)
        cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        r = cp.add_run(); r.text = caption
        set_font(r, 9, color=GRAY)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def make_grid():
    """5x3 contact sheet of the 15 source photos for slide 2."""
    GRID.parent.mkdir(parents=True, exist_ok=True)
    paths = sorted(IMAGES.glob("*.jpg"))[:15]
    cols, rows = 5, 3
    cw, ch, pad = 360, 240, 8
    sheet = Image.new("RGB", (cols * cw + pad * (cols + 1),
                              rows * ch + pad * (rows + 1)), (245, 247, 250))
    for i, p in enumerate(paths):
        with Image.open(p) as im:
            im.draft("RGB", (cw * 2, ch * 2))
            im = im.convert("RGB")
            iw, ih = im.size
            s = min(cw / iw, ch / ih)
            im = im.resize((max(1, int(iw * s)), max(1, int(ih * s))), Image.LANCZOS)
        r, c = divmod(i, cols)
        ox = pad + c * (cw + pad) + (cw - im.width) // 2
        oy = pad + r * (ch + pad) + (ch - im.height) // 2
        sheet.paste(im, (ox, oy))
    sheet.save(GRID, quality=88)
    return GRID


def main():
    prs = Presentation()
    prs.slide_width = Emu(int(SW * EMU_IN))
    prs.slide_height = Emu(int(SH * EMU_IN))

    # ---- Slide 1: title ----
    s = blank(prs)
    rect(s, 0, 0, SW, SH, WHITE)
    rect(s, 0, 2.55, SW, 0.05, RED)
    rect(s, 0, 4.30, SW, 0.02, BLUE)
    tf = textbox(s, 1.0, 2.7, SW - 2.0, 1.5, MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "夜間運動模糊的人臉修復"
    set_font(r, 40, bold=True, color=DARK)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(6)
    r = p2.add_run(); r.text = "回歸式去模糊結合生成式先驗"
    set_font(r, 20, color=BLUE)
    meta = textbox(s, 1.0, 4.5, SW - 2.0, 1.3, MSO_ANCHOR.TOP)
    mp = meta.paragraphs[0]; mp.alignment = PP_ALIGN.CENTER
    r = mp.add_run(); r.text = "113950011 鄭名翔　·　[id] [teammate]　·　[id] [teammate]"
    set_font(r, 15, color=DARK)
    mp2 = meta.add_paragraph(); mp2.alignment = PP_ALIGN.CENTER; mp2.space_before = Pt(8)
    r = mp2.add_run(); r.text = "NYCU 影像處理 · Term Project"
    set_font(r, 13, color=GRAY)
    notes(s, "0:00–0:40 題目與目標、作業範例就是臉。")

    # ---- Slide 2: problem & goal ----
    s = blank(prs)
    title_bar(s, "問題與目標", 2)
    bullets(s, [
        "15 張夜間真實照片，6K–8K，含 motion blur，無 ground truth（不能用 PSNR/SSIM）",
        "作業範例正是 blurry face → sharp face，我們以此為目標",
        "目標：主體可信還原的前提下，做出前後差異最明顯的 before/after",
        ("挑模糊最重、但仍可修復的主體", 1),
    ], 0.6, 1.5, 5.7, 4.5, size=16)
    fit_image(s, make_grid(), 6.5, 1.55, 6.4, 4.7, caption="題目提供的 15 張夜間模糊照片")
    notes(s, "0:00–0:40 題目與目標、作業範例就是臉。")

    # ---- Slide 3: baseline & ceiling ----
    s = blank(prs)
    title_bar(s, "基礎方法與它的天花板", 3)
    bullets(s, [
        "成功執行 FFTformer（CVPR 2023，頻域 transformer，RealBlur-J 權重）",
        "也比較過 DarkIR（只提亮）、MISCFilter（去模糊有限）",
        "回歸式學 blur→sharp；臉的高頻被抹除後，輸出偏軟、帶噪（接近上限）",
    ], 0.6, 1.5, 12.1, 1.7, size=16)
    fit_image(s, FIG / "face_fft_vs_supir.jpg", 1.4, 3.35, 10.5, 3.0,
              caption="RAW ｜ FFTformer ｜ FFTformer→SUPIR")
    notes(s, "0:40–1:30 FFTformer 與回歸天花板：高頻被抹除後回歸模型補不出，輸出偏軟帶噪。")

    # ---- Slide 4: method ----
    s = blank(prs)
    title_bar(s, "方法：回歸去模糊 + 生成式臉部先驗", 4)
    bullets(s, [
        "裁臉 → FFTformer（去模糊、給乾淨結構）→ SUPIR-v0F（擴散先驗補臉部高頻，高 control 貼合五官）→ 羽化合成 → 輕度調色",
        "兩階段互補：先把 kernel 縮回範圍給乾淨結構，再補細節；high control 控制幻覺",
    ], 0.6, 1.5, 12.1, 1.5, size=16)
    fit_image(s, FIG / "face_pipeline.png", 0.6, 3.2, 12.1, 3.2)
    notes(s, "1:30–2:30 方法（兩階段）+ 兩張結果：先縮 kernel 給乾淨結構，再用擴散先驗補高頻，high control 緊錨原五官。")

    # ---- Slide 5: results ----
    s = blank(prs)
    title_bar(s, "結果", 5)
    bullets(s, [
        "繳交 2 張不同場景的「模糊路人 → 清晰真人臉」",
        "#2（女子）上半臉為生成，下一頁誠實揭露",
    ], 0.6, 1.5, 12.1, 1.1, size=16)
    fit_image(s, FIG / "face_result_man.jpg", 0.6, 2.75, 6.0, 3.7, caption="#1 夜市男子")
    fit_image(s, FIG / "face_result_woman.jpg", 6.9, 2.75, 5.9, 3.7, caption="#2 霓虹巷弄女子")
    notes(s, "1:30–2:30 兩張結果：模糊路人浮現成清晰真人臉，背景保留動態模糊。")

    # ---- Slide 6: fidelity (key) ----
    s = blank(prs)
    title_bar(s, "誠實面：真實修復 vs 生成合成（重點）", 6)
    bullets(s, [
        "臉部細節是擴散先驗合成（prior-guided，與範例 GFPGAN/GPEN 同類），非解卷積",
        "#1 男子：五官 layout 真實、高 control 緊錨；高頻仍為先驗合成",
        "#2 女子：眼睛 / 上半臉為生成（原圖被運動重影破壞）",
        "重點：高 NR-IQA 不代表真實還原",
    ], 0.6, 1.5, 6.6, 4.5, size=15)
    fit_image(s, FIG / "face_seed_consistency.jpg", 7.4, 2.1, 5.5, 3.6,
              caption="不同 seed：受結構約束處大致一致（一致 ≠ 忠實）")
    notes(s, "2:30–3:40 誠實面：務必講清楚 #2 女子的眼睛是生成的；一致不等於忠實，高 NR-IQA 不代表真實還原。")

    # ---- Slide 7: engineering difficulties (key) ----
    s = blank(prs)
    title_bar(s, "解決的實作困難（重點）", 7)
    bullets(s, [
        "16 GB 顯存跑 6–8K：FFTformer overlap-blend tiling；SUPIR tiled VAE / tiled sampling",
        "Blackwell sm_120：需 PyTorch 2.11+cu128；ComfyUI/SUPIR 與 FFTformer 相依衝突",
        ("→ 兩個隔離的 conda env（deblur / comfy）", 1),
        "生成式幻覺控制：裁臉（背景不亂修）+ 高 control + 裁掉幻覺邊緣",
        ("整圖會在主體旁生出第二張臉，故改裁臉局部修復", 1),
    ], 0.6, 1.6, 12.1, 4.6, size=17)
    notes(s, "2:30–3:40 解決的困難：顯存 tiling、Blackwell 雙環境、幻覺控制——這是課堂加分重點。")

    # ---- Slide 8: rejected directions ----
    s = blank(prs)
    title_bar(s, "驗證過、但沒採用的方向", 8)
    bullets(s, [
        "純車輛去模糊（05、08）：真實去模糊但對比太小",
        "玻璃反射（15）：layer separation 非模糊，去不掉",
        "整張 SUPIR：把背景人群幻覺成扭曲臉",
        "重點：說明什麼有效、什麼無效、為什麼",
    ], 0.6, 1.6, 7.2, 4.5, size=16)
    fit_image(s, FIG / "face_rejected.jpg", 8.0, 1.5, 4.9, 4.9,
              caption="玻璃反射（左）｜車輛去模糊對比小（右）")
    notes(s, "3:40–4:30 驗證過但沒採用：對比太小 / layer separation 去不掉 / 整圖幻覺——說明取捨。")

    # ---- Slide 9: conclusion ----
    s = blank(prs)
    title_bar(s, "結論與限制", 9)
    bullets(s, [
        "貢獻：回歸→生成 兩階段人臉修復 + 誠實的 fidelity 界線",
        "限制：玻璃反射需 layer separation；被重影破壞的五官無法真實還原，只能可信生成",
        "未來：臉部專用修復（fidelity 旋鈕）、真實 paired night-blur 微調、reflection removal",
    ], 0.6, 1.7, 12.1, 4.3, size=18)
    rect(s, 0.6, 6.0, 12.1, 0.03, RED)
    end = textbox(s, 0.6, 6.1, 12.1, 0.6, MSO_ANCHOR.MIDDLE)
    ep = end.paragraphs[0]; ep.alignment = PP_ALIGN.CENTER
    r = ep.add_run(); r.text = "謝謝聆聽"
    set_font(r, 16, bold=True, color=BLUE)
    notes(s, "4:30–5:00 結論與限制：兩階段貢獻 + 誠實 fidelity 界線；玻璃反射與被重影抹除的五官是已知限制。")

    n = len(prs.slides._sldIdLst)
    prs.save(OUT)
    print(f"[ok] {OUT}  ({OUT.stat().st_size/1024:.0f} KB, {n} slides)")


if __name__ == "__main__":
    main()
